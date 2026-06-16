"""Перекраивает КudaEdem_presentation.pptx под агент-ориентированный нарратив, СОХРАНЯЯ дизайн
(оформление зашито в png-фоны каждого слайда, поэтому работаем внутри их файла, а не строим заново).

    python agent/scripts/restructure_presentation.py
        ->  KudaEdem_agent.pptx   (их дизайн, ~12 слайдов, без «КРИТЕРИЙ N»)
"""
from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation

SRC = Path("KudaEdem_presentation.pptx")
OUT = Path("KudaEdem_agent.pptx")
R_ID = "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"

# Желаемый порядок по 1-based номерам исходных слайдов (см. карту в чате).
KEEP_1BASED = [1, 2, 5, 3, 4, 8, 7, 9, 10, 15, 16, 17]

prs = Presentation(str(SRC))
sld_id_lst = prs.slides._sldIdLst
sld_ids = list(sld_id_lst)  # текущий порядок, 0-indexed = номер слайда - 1
keep_idx = [n - 1 for n in KEEP_1BASED]
keep_set = set(keep_idx)

# 1) удалить ненужные слайды (снять relationship → часть не попадёт в файл, медиа подрежется).
for i, sld_id in enumerate(sld_ids):
    if i not in keep_set:
        prs.part.drop_rel(sld_id.get(R_ID))
        sld_id_lst.remove(sld_id)

# 2) переставить оставшиеся в нужном порядке (remove+append уже присутствующих элементов).
for i in keep_idx:
    el = sld_ids[i]
    sld_id_lst.remove(el)
    sld_id_lst.append(el)

# 3) убрать префикс «КРИТЕРИЙ N · » в кикерах, оставив название секции.
prefix = re.compile(r"^\s*КРИТЕРИЙ[^·]*·\s*")
for slide in prs.slides:
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if run.text.strip().startswith("КРИТЕРИЙ"):
                    run.text = prefix.sub("", run.text)

prs.save(str(OUT))
print(f"saved {OUT} · slides={len(prs.slides._sldIdLst)}")
