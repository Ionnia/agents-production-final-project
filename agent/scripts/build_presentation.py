"""Презентация КудаЕдем к защите — с нуля, фокус на агенте, острый текст без воды.

    python agent/scripts/build_presentation.py   ->  KudaEdem_agent.pptx

Чистый светлый дизайн (их png-фоны из кода не воспроизвести). Каждый слайд — одна мысль,
конкретика и цифры, минимум текста.
"""
from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

INK = RGBColor(0x1B, 0x23, 0x33)
MUTED = RGBColor(0x66, 0x71, 0x82)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT = RGBColor(0x3B, 0x5B, 0xDB)
GOOD = RGBColor(0x0C, 0xA6, 0x78)
WARN = RGBColor(0xE8, 0x92, 0x0C)
BORDER = RGBColor(0xDF, 0xE5, 0xEF)
LLM_F, LLM_L = RGBColor(0xE9, 0xED, 0xFF), ACCENT
DET_F, DET_L = RGBColor(0xFF, 0xF3, 0xDE), WARN
DATA_F, DATA_L = RGBColor(0xE2, 0xF7, 0xEF), GOOD
NEUT_F = RGBColor(0xF4, 0xF6, 0xFB)

W, H = Inches(13.333), Inches(7.5)
prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]
_page = 0


def slide(section=""):
    global _page
    _page += 1
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = WHITE
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.16), H)
    bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT; bar.line.fill.background(); bar.shadow.inherit = False
    _txt(s, Inches(0.7), Inches(7.0), Inches(8), Inches(0.4), [("КудаЕдем", 11, MUTED, False)])
    _txt(s, Inches(12.4), Inches(7.0), Inches(0.7), Inches(0.4), [(str(_page), 11, MUTED, False)], align=PP_ALIGN.RIGHT)
    return s


def _txt(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=6, leading=None):
    tb = s.shapes.add_textbox(l, t, w, h); tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, (txt, size, color, bold) in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(space)
        if leading:
            p.line_spacing = leading
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.name = "Calibri"
    return tb


def kicker(s, txt):
    _txt(s, Inches(0.75), Inches(0.6), Inches(11.8), Inches(0.4), [(txt.upper(), 13, ACCENT, True)])


def headline(s, txt, size=36, top=1.05):
    _txt(s, Inches(0.75), Inches(top), Inches(12.0), Inches(1.5), [(txt, size, INK, True)], leading=1.0)


def body(s, items, top=2.7, size=19, width=11.6, gap=12):
    runs = [("•  " + t, size, c, b) for (t, c, b) in items]
    _txt(s, Inches(0.9), Inches(top), Inches(width), Inches(4), runs, space=gap, leading=1.05)


def metric(s, val, color=GOOD, top=5.7):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(top), Inches(4.6), Inches(0.95))
    b.fill.solid(); b.fill.fore_color.rgb = WHITE; b.line.color.rgb = color; b.line.width = Pt(1.5); b.shadow.inherit = False
    tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    for seg, col in [("outcome ", MUTED), (val.split("/")[0].strip().split()[-1], color),
                     ("   entity ", MUTED), (val.split("/")[1].strip(), color)]:
        r = p.add_run(); r.text = seg; r.font.size = Pt(22); r.font.bold = True; r.font.color.rgb = col


def chip(s, x, y, w, h, head, sub, fill, line):
    b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    b.fill.solid(); b.fill.fore_color.rgb = fill; b.line.color.rgb = line; b.line.width = Pt(1.25); b.shadow.inherit = False
    tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = head; r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = INK
    if sub:
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run(); r2.text = sub; r2.font.size = Pt(11); r2.font.color.rgb = MUTED


def flow(s, steps, top, bw=1.62, gap=0.16, x0=0.75, h=1.45):
    bw, gap, x0 = Inches(bw), Inches(gap), Inches(x0); y = Inches(top)
    for i, (head, sub, fill, line) in enumerate(steps):
        x = Emu(int(x0) + i * (int(bw) + int(gap)))
        chip(s, x, y, bw, Inches(h), head, sub, fill, line)
        if i < len(steps) - 1:
            _txt(s, Emu(int(x) + int(bw)), y, gap, Inches(h), [("→", 16, ACCENT, True)],
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


# ── 1. Титул ─────────────────────────────────────────────────────────────────────────────────
s = slide(); kicker(s, "Промышленная разработка агентов · защита")
_txt(s, Inches(0.75), Inches(2.3), Inches(12), Inches(1.6), [("КудаЕдем", 70, INK, True)])
_txt(s, Inches(0.78), Inches(3.9), Inches(11.6), Inches(1.6),
     [("ИИ-агент, который планирует поездку и сам решает:", 26, INK, False),
      ("рекомендовать · уточнить · отказать · эскалировать.", 26, ACCENT, True)], space=4, leading=1.05)
_txt(s, Inches(0.78), Inches(5.8), Inches(8), Inches(0.6), [("Итог: outcome 0.90 · entity 1.0", 18, MUTED, False)])

# ── 2. Проблема ──────────────────────────────────────────────────────────────────────────────
s = slide(); kicker(s, "Проблема")
headline(s, "Подбор поездки — это решение, а не поиск")
body(s, [
    ("Требования конфликтуют: «только прямой + только 5★ + бюджет 200к» одновременно невыполнимо.", INK, False),
    ("В свободном чате данных не хватает; иногда задача в принципе невыполнима.", INK, False),
    ("Агент должен выбрать ТИП ответа — рекомендовать / уточнить / отказать / эскалировать,", INK, False),
    ("а не просто «найти, что дешевле».", ACCENT, True),
], top=2.6)

# ── 3. Система ───────────────────────────────────────────────────────────────────────────────
s = slide(); kicker(s, "Система · одним слайдом")
headline(s, "Три сервиса — в LLM ходит только агент")
flow(s, [("Frontend", "чат · карта", NEUT_F, BORDER), ("Backend (BFF)", "данные · планы · доступ", NEUT_F, BORDER),
         ("Agent Service", "граф · LLM · RAG", LLM_F, LLM_L)], top=2.9, bw=3.4, gap=0.45, x0=0.85, h=1.5)
body(s, [
    ("Фронт говорит только с бэкендом. Бэкенд владеет данными и драйвит агента.", MUTED, False),
    ("Контракты между сервисами зафиксированы в OpenAPI → команды работали параллельно. Дальше — только агент.", MUTED, False),
], top=4.9, size=18)

# ── 4. Метрики + подход ──────────────────────────────────────────────────────────────────────
s = slide(); kicker(s, "Как мерим качество")
headline(s, "Две метрики — и они тянут в разные стороны")
body(s, [
    ("outcome — верный ТИП ответа (рекомендация / уточнение / отказ / эскалация).", INK, False),
    ("entity — верные ID: тот самый рейс и отель, что в эталоне.", INK, False),
], top=2.5)
flow(s, [("B1", "single", NEUT_F, BORDER), ("B2", "+validate", NEUT_F, BORDER),
         ("B3", "мультиагент", NEUT_F, BORDER), ("Final", "синтез", DATA_F, DATA_L)],
     top=4.2, bw=2.5, gap=0.35, x0=0.85, h=1.2)
_txt(s, Inches(0.9), Inches(5.9), Inches(11.5), Inches(0.8),
     [("4 бейзлайна, один eval-набор (20 кейсов). Задача — поднять ОБЕ метрики. Оказалось, это разные вещи.", 17, MUTED, False)])

# ── 5. B1 ────────────────────────────────────────────────────────────────────────────────────
s = slide(); kicker(s, "Бейзлайн B1")
headline(s, "Один агент с инструментами — и провал на деньгах")
body(s, [
    ("ReAct-агент сам зовёт поиск офферов и RAG-правила, выдаёт план.", INK, False),
    ("Провал: доверили LLM считать бюджет — он выдумывал перерасход и ложно отказывал.", WARN, False),
], top=2.7)
metric(s, "0.85 / 0.67", WARN)

# ── 6. B2 ────────────────────────────────────────────────────────────────────────────────────
s = slide(); kicker(s, "Бейзлайн B2")
headline(s, "Бюджет считает код, а не модель")
body(s, [
    ("draft (LLM) → validate (детерминированно: бюджет, ограничения) → переиграть.", INK, False),
    ("Инсайт: арифметику нельзя отдавать LLM. outcome 0.85 → 0.95.", GOOD, False),
    ("Но entity всё ещё 0.67 — одиночный агент путает, какой именно отель/рейс.", WARN, False),
], top=2.7)
metric(s, "0.95 / 0.67", GOOD)

# ── 7. B3 ────────────────────────────────────────────────────────────────────────────────────
s = slide(); kicker(s, "Бейзлайн B3")
headline(s, "По агенту на перелёт, отель и тур")
body(s, [
    ("Узкие специалисты подбирают точнее → entity взлетел до 1.0.", GOOD, False),
    ("Но outcome упал до ~0.75: решение принимается по текстовым сводкам, далеко от данных → шумит.", WARN, False),
], top=2.7)
metric(s, "~0.75 / 1.0", WARN)

# ── 8. Final ─────────────────────────────────────────────────────────────────────────────────
s = slide(); kicker(s, "Финал")
headline(s, "Соединил сильные стороны B2 и B3")
body(s, [
    ("entity — от специалистов: они находят правильные id.", INK, False),
    ("outcome — от детерминированных ФАКТОВ (feasibility: стоимость, бюджет, нарушения).", INK, False),
    ("Агент решает тип ответа по этим фактам, а не по догадке.", INK, False),
], top=2.7)
metric(s, "0.90 / 1.0", GOOD)

# ── 9. Принцип ───────────────────────────────────────────────────────────────────────────────
s = slide(); kicker(s, "Главная идея")
_txt(s, Inches(0.75), Inches(2.5), Inches(12), Inches(2.2),
     [("Факты считает код.", 50, INK, True), ("Решение принимает агент.", 50, ACCENT, True)], space=6, leading=1.05)
_txt(s, Inches(0.78), Inches(5.0), Inches(11.6), Inches(1.6),
     [("Арифметику и проверки — детерминированный калькулятор. Тип ответа и подбор — модель, "
       "заземлённая на эти факты.", 20, MUTED, False),
      ("Это сняло размен outcome ↔ entity: оба показателя стали высокими без подгонки под тест.", 20, MUTED, False)],
     space=6, leading=1.05)

# ── 10. Архитектура ──────────────────────────────────────────────────────────────────────────
s = slide(); kicker(s, "Архитектура агента")
headline(s, "Один граф (LangGraph)")
flow(s, [("resolve", "понять запрос", DET_F, DET_L), ("context", "данные + RAG", DATA_F, DATA_L),
         ("intent", "бюджет, риск", LLM_F, LLM_L), ("специалисты", "подбор офферов", LLM_F, LLM_L),
         ("feasibility", "ФАКТЫ", DET_F, DET_L), ("supervisor", "+ critic", LLM_F, LLM_L),
         ("исход", "ответ", NEUT_F, BORDER)], top=2.55, bw=1.6, gap=0.12, x0=0.55, h=1.5)
body(s, [
    ("Жёлтое — детерминированный код (факты), синее — LLM. Зелёное — данные/RAG.", MUTED, False),
    ("supervisor выбирает тип ответа ПО ФАКТАМ; critic сверяет и отправляет на переплан (×2).", MUTED, False),
    ("Слева быстрые ветки: вопрос о правилах · уточнение направления вариантами · болталка.", MUTED, False),
], top=4.45, size=17, gap=8)

# ── 11. Метрики ──────────────────────────────────────────────────────────────────────────────
s = slide(); kicker(s, "Результаты")
headline(s, "Что дал каждый шаг")
data = [
    ("Вариант", "outcome", "entity", "идея шага"),
    ("B1 — один агент + RAG", "0.85", "0.67", "базовый агент"),
    ("B2 — + детерм. validate", "0.95", "0.67", "факты → стабильный outcome"),
    ("B3 — мультиагент", "~0.75", "1.0", "специалисты → entity"),
    ("Final — синтез", "0.90", "1.0", "факты + специалисты"),
]
tbl = s.shapes.add_table(len(data), 4, Inches(0.75), Inches(2.4), Inches(11.8), Inches(3.0)).table
tbl.columns[0].width = Inches(4.2); tbl.columns[1].width = Inches(1.7)
tbl.columns[2].width = Inches(1.7); tbl.columns[3].width = Inches(4.2)
for ri, row in enumerate(data):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if ri == 0 else (DATA_F if ri == len(data) - 1 else (NEUT_F if ri % 2 else WHITE))
        p = cell.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if ci in (0, 3) else PP_ALIGN.CENTER
        r = p.add_run(); r.text = val
        r.font.size = Pt(15); r.font.bold = (ri == 0 or ri == len(data) - 1)
        r.font.color.rgb = WHITE if ri == 0 else (GOOD if ci in (1, 2) and ri >= 1 else INK)
_txt(s, Inches(0.78), Inches(5.7), Inches(11.7), Inches(1.0),
     [("B2 — лучший outcome, но слабый entity. B3 — наоборот. Final снял размен: 0.90 / 1.0.", 18, INK, True),
      ("Цифры — про размен метрик, не строгий A/B (QA-набор дорабатывался).", 14, MUTED, False)], space=6)

# ── 12. Честно + итог ────────────────────────────────────────────────────────────────────────
s = slide(); kicker(s, "Честно · итог")
headline(s, "Что осталось — и главный вывод")
body(s, [
    ("Остаточные промахи: граница «отказать ↔ уточнить» спорна в самой разметке.", INK, False),
    ("LLM-вариативность на пограничных кейсах — мерим честно, как диапазон.", INK, False),
    ("Дальше: персистентная память графа, валидация плана через реальный Contract B.", MUTED, False),
], top=2.5, gap=10)
b = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(5.0), Inches(11.8), Inches(1.5))
b.fill.solid(); b.fill.fore_color.rgb = DATA_F; b.line.color.rgb = GOOD; b.line.width = Pt(1.5); b.shadow.inherit = False
tf = b.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE; tf.word_wrap = True
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
r = p.add_run(); r.text = "Один принцип — факты в коде, решение у агента — дал и outcome, и entity."
r.font.size = Pt(24); r.font.bold = True; r.font.color.rgb = INK

out = Path("KudaEdem_agent.pptx")
prs.save(out)
print(f"saved {out} · slides={len(prs.slides._sldIdLst)}")
